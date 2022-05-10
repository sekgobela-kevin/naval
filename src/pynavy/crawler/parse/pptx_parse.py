from io import StringIO

import pptx

from .parse_base import Parse_Base
from ..fetch.master_fetch import Master_Fetch

class PPTX_Parse(Parse_Base):
    '''Parses html data from fetch object'''
    def __init__(self, fetch_obj) -> None:
        super().__init__(fetch_obj)

    def create_doc(self, *args, **kwarg):
        '''Return object to use when parsing fetch object contents.'''
        return pptx.Presentation(self.fetch_obj.get_file())

    def get_text(self) -> str or bytes:
        '''Retuns text version of fetch object'''
        text_list = []
        for slide in self.doc.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_list.append(shape.text)
        return "\n".join(text_list)

if __name__ == "__main__":
    from ..fetch.web_fetch import Web_Fetch

    # create fetch object
    
    url = "/mnt/c/Users/sekgobela kevin/Downloads/Lecture 1, What is History of Education (1).pptx"
    fetch_obj = Master_Fetch.get_fetch_object(url)
    # this performs request for data
    fetch_obj.request()
    print(len(fetch_obj.fetch()))

    # create parse object from fetch object
    parse_obj = PPTX_Parse(fetch_obj)
    print(parse_obj.get_text())